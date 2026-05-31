from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BLUE = RGBColor(64, 112, 196)
GREEN = RGBColor(78, 156, 96)
ORANGE = RGBColor(239, 155, 53)
PURPLE = RGBColor(124, 98, 204)

BLUE_FILL = RGBColor(245, 249, 255)
GREEN_FILL = RGBColor(247, 252, 247)
ORANGE_FILL = RGBColor(255, 250, 243)
PURPLE_FILL = RGBColor(249, 247, 255)

TEXT_BLUE = RGBColor(46, 86, 162)
TEXT_GREEN = RGBColor(66, 138, 84)
TEXT_ORANGE = RGBColor(220, 140, 44)
TEXT_PURPLE = RGBColor(110, 86, 187)
CARD_FILL = RGBColor(255, 255, 255)


def add_round_box(slide, x, y, w, h, fill, line, radius=0.06, line_width=1.0):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_width)
    shape.adjustments[0] = radius
    return shape


def set_shape_text(shape, text, size, color, bold=False, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_connector(slide, x1, y1, x2, y2, color, width=1.4):
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    shape.line.end_arrowhead = True
    return shape


def add_layer(slide, y, border, fill, left_text, text_color, right_labels):
    outer_x = Inches(0.45)
    outer_w = Inches(12.45)
    outer_h = Inches(1.08)
    left_w = Inches(1.85)
    gap = Inches(0.36)
    card_h = Inches(0.62)
    card_y = y + Inches(0.23)
    card_x = outer_x + left_w + Inches(0.35)
    card_w = Inches(2.14)

    add_round_box(slide, outer_x, y, outer_w, outer_h, fill, border, radius=0.04, line_width=1.0)
    left = add_round_box(slide, outer_x, y, left_w, outer_h, fill, border, radius=0.04, line_width=1.0)
    set_shape_text(left, left_text, 17, text_color, True)

    cards = []
    for i, label in enumerate(right_labels):
        x = card_x + i * (card_w + gap)
        card = add_round_box(slide, x, card_y, card_w, card_h, CARD_FILL, border, radius=0.035, line_width=0.9)
        set_shape_text(card, label, 13.5, text_color, True)
        cards.append(card)
    return cards


def row_center(cards):
    first = cards[0]
    last = cards[-1]
    return (first.left + (last.left + last.width)) / 2


def build_deck(output_path: str | Path = "output/arch_like_original_editable.pptx") -> Path:
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(255, 255, 255)

    layer1_y = Inches(0.32)
    layer2_y = Inches(2.10)
    layer3_y = Inches(3.88)
    layer4_y = Inches(5.66)

    layer1 = add_layer(slide, layer1_y, BLUE, BLUE_FILL, "业务层", TEXT_BLUE, ["客户门户", "移动应用", "运营驾驶舱"])
    layer2 = add_layer(slide, layer2_y, GREEN, GREEN_FILL, "应用层", TEXT_GREEN, ["统一门户", "服务编排", "监控中心", "工单系统"])
    layer3 = add_layer(slide, layer3_y, ORANGE, ORANGE_FILL, "能力层", TEXT_ORANGE, ["API 网关", "身份认证", "消息中台", "分析引擎"])
    layer4 = add_layer(slide, layer4_y, PURPLE, PURPLE_FILL, "数据层", TEXT_PURPLE, ["关系数据库", "缓存", "数据仓库", "对象存储"])

    add_connector(slide, row_center(layer1), layer1_y + Inches(1.08), row_center(layer1), layer2_y, BLUE, 1.2)
    add_connector(slide, row_center(layer2), layer2_y + Inches(1.08), row_center(layer2), layer3_y, GREEN, 1.2)
    add_connector(slide, row_center(layer3), layer3_y + Inches(1.08), row_center(layer3), layer4_y, ORANGE, 1.2)

    prs.save(output)
    return output


if __name__ == "__main__":
    out = build_deck()
    print(out)
