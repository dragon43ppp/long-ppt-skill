from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


W = Inches(20.0)
H = Inches(8.6)

NAVY = RGBColor(0, 81, 130)
DEEP = RGBColor(0, 96, 150)
BLUE = RGBColor(26, 167, 209)
BLUE2 = RGBColor(187, 238, 247)
GREEN = RGBColor(113, 184, 54)
GREEN2 = RGBColor(219, 244, 207)
PANEL_BLUE = RGBColor(246, 253, 255)
PANEL_GREEN = RGBColor(249, 255, 248)
LINE_BLUE = RGBColor(139, 204, 227)
LINE_GREEN = RGBColor(139, 210, 184)
TEXT = RGBColor(35, 76, 100)
WHITE = RGBColor(255, 255, 255)


def set_text(shape, text, size=10, color=TEXT, bold=False, align=None):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    lines = str(text).split("\n")
    p = tf.paragraphs[0]
    p.text = lines[0]
    if align is not None:
        p.alignment = align
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        if align is not None:
            p.alignment = align
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = "Microsoft YaHei"
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color


def textbox(slide, x, y, w, h, text, size=10, color=TEXT, bold=False, align=None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(shape, text, size, color, bold, align)
    return shape


def box(slide, x, y, w, h, fill=WHITE, line=LINE_BLUE, radius=True, width=1.0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(width)
    return shape


def line(slide, x1, y1, x2, y2, color=DEEP, width=1.4):
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    return shape


def section_title(slide, y, text):
    line(slide, 0.12, y + 0.08, 9.15, y + 0.08, DEEP, 2.0)
    line(slide, 10.9, y + 0.08, 19.88, y + 0.08, DEEP, 2.0)
    textbox(slide, 9.45, y - 0.03, 1.25, 0.24, text, 12, NAVY, True, PP_ALIGN.CENTER)


def tag(slide, x, y, w, text, fill, line_color=None):
    shape = box(slide, x, y, w, 0.24, fill, line_color or fill, True, 0.8)
    set_text(shape, text, 6.8, NAVY if fill != GREEN else WHITE, True, PP_ALIGN.CENTER)
    return shape


def content_card(slide, x, y, w, h, title, body, tags, accent=LINE_BLUE, tag_fill=BLUE2):
    box(slide, x, y, w, h, WHITE, accent, True, 1.0)
    textbox(slide, x + 0.12, y + 0.1, w - 0.24, 0.28, title, 11.3, NAVY, True)
    textbox(slide, x + 0.12, y + 0.43, w - 0.24, h - 0.77, body, 7.7, TEXT, False)
    if tags:
        gap = 0.08
        tw = (w - 0.28 - gap * (len(tags) - 1)) / len(tags)
        for i, item in enumerate(tags):
            tag(slide, x + 0.14 + i * (tw + gap), y + h - 0.33, tw, item, tag_fill, accent)


def phase_panel(slide, x, y, w, h, title, subtitle, badge, badge_color, fill, line_color):
    box(slide, x, y, w, h, fill, line_color, True, 1.2)
    textbox(slide, x + 0.16, y + 0.22, w * 0.42, 0.34, title, 15, NAVY, True)
    textbox(slide, x + 0.16, y + 0.68, w - 0.32, 0.54, subtitle, 7.8, TEXT)
    badge_shape = box(slide, x + w - 1.0, y + 0.18, 0.78, 0.30, badge_color, badge_color, True, 0.8)
    set_text(badge_shape, badge, 6.2, WHITE, True, PP_ALIGN.CENTER)


def build_deck(output_path: str | Path = "output/attached_architecture_editable.pptx") -> Path:
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    box(slide, 0.08, 0.08, 19.84, 5.38, WHITE, DEEP, True, 1.4)
    section_title(slide, 0.22, "阶段规划")

    phase_panel(
        slide, 0.28, 0.56, 6.25, 4.68,
        "一期建设",
        "围绕统一纳管、实时监测、主动预警和高效排障，先把基础可视、可证、可定界能力做实。",
        "优先落地", GREEN, PANEL_GREEN, LINE_GREEN,
    )
    phase_panel(
        slide, 6.80, 0.56, 6.25, 4.68,
        "二期扩展",
        "逐步扩展到更多业务系统、跨域资源与专题分析能力，形成更完整的平台化支撑视图。",
        "平台扩展", BLUE, PANEL_BLUE, LINE_BLUE,
    )
    phase_panel(
        slide, 13.32, 0.56, 6.25, 4.68,
        "三期演进",
        "在稳定的数据底座和流程协同之上，进一步引入智能分析、自动化闭环和知识沉淀能力。",
        "智能演进", DEEP, PANEL_BLUE, LINE_BLUE,
    )

    content_card(
        slide, 0.45, 1.82, 2.86, 1.38,
        "监测覆盖",
        "优先覆盖关键链路、核心业务和重点区域，建立统一采集与统一展示能力。",
        ["链路", "业务", "重点区域"],
        LINE_GREEN,
        GREEN2,
    )
    content_card(
        slide, 3.42, 1.82, 2.86, 1.38,
        "预警能力",
        "建立实时告警、阈值联动和异常识别能力，为后续闭环处置打基础。",
        ["实时", "阈值", "异常"],
        LINE_GREEN,
        GREEN2,
    )
    content_card(
        slide, 6.98, 1.82, 2.86, 1.38,
        "专题分析",
        "围绕业务连续性、跨域定位和根因排查，建设专题化分析视图。",
        ["业务", "跨域", "根因"],
    )
    content_card(
        slide, 9.95, 1.82, 2.86, 1.38,
        "平台协同",
        "形成统一对象视图、统一证据链和统一协同入口，支撑多角色联动。",
        ["对象", "证据", "协同"],
    )
    content_card(
        slide, 13.52, 1.82, 2.86, 1.38,
        "智能闭环",
        "在稳定运行基础上引入知识沉淀、建议输出和流程闭环优化能力。",
        ["知识", "建议", "闭环"],
    )
    content_card(
        slide, 16.49, 1.82, 2.86, 1.38,
        "运营支撑",
        "面向管理侧输出专题看板、汇报材料和运营评价依据。",
        ["看板", "汇报", "评价"],
    )

    section_title(slide, 5.72, "能力结构")
    box(slide, 0.38, 6.00, 19.20, 2.10, RGBColor(250, 253, 255), RGBColor(171, 214, 234), True, 1.0)

    content_card(slide, 0.62, 6.26, 3.10, 1.48, "统一采集", "多源接入、统一对象、统一标签，支撑后续分析与展示。", ["采集", "对象", "标签"])
    content_card(slide, 4.02, 6.26, 3.10, 1.48, "统一分析", "指标、链路、日志与事件联动，形成问题证据链。", ["指标", "链路", "日志"])
    content_card(slide, 7.42, 6.26, 3.10, 1.48, "统一协同", "告警、工单、复盘和知识沉淀串成闭环。", ["告警", "工单", "复盘"])
    content_card(slide, 10.82, 6.26, 3.10, 1.48, "统一展示", "面向运维、管理和业务角色输出差异化视图。", ["运维", "管理", "业务"])
    content_card(slide, 14.22, 6.26, 3.10, 1.48, "统一运营", "支撑阶段汇报、运营评价和持续优化。", ["汇报", "评价", "优化"])

    prs.save(output)
    return output


if __name__ == "__main__":
    out = build_deck()
    print(out)
