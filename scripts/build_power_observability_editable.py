from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


W = Inches(16.0)
H = Inches(8.92)

NAVY = RGBColor(6, 54, 143)
BLUE = RGBColor(22, 102, 224)
BLUE2 = RGBColor(35, 132, 235)
LIGHT_BLUE = RGBColor(235, 246, 255)
LINE = RGBColor(199, 224, 250)
TEXT = RGBColor(42, 55, 75)
WHITE = RGBColor(255, 255, 255)


def set_text(shape, text, size=10, color=TEXT, bold=False, align=None):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    lines = str(text).split("\n")
    p = tf.paragraphs[0]
    p.text = lines[0]
    if align is not None:
        p.alignment = align
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        if align is not None:
            p.alignment = align
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = "Microsoft YaHei"
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color


def textbox(slide, x, y, w, h, text, size=10, color=TEXT, bold=False, align=None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(shape, text, size, color, bold, align)
    return shape


def box(slide, x, y, w, h, fill=WHITE, line=LINE, radius=True, width=1.0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(width)
    return shape


def connector(slide, x1, y1, x2, y2, color=BLUE2, width=1.2):
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    return shape


def icon_circle(slide, x, y, label, size=0.8):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.color.rgb = BLUE
    set_text(shape, label, 18, WHITE, True, PP_ALIGN.CENTER)
    return shape


def risk_card(slide, x, y, no, title, body, icon):
    box(slide, x, y, 3.48, 1.38, WHITE, LINE, True, 1.0)
    icon_circle(slide, x + 0.15, y + 0.26, icon, 0.75)
    n = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.18), Inches(y + 0.26), Inches(0.28), Inches(0.28))
    n.fill.solid()
    n.fill.fore_color.rgb = NAVY
    n.line.color.rgb = NAVY
    set_text(n, str(no), 9, WHITE, True, PP_ALIGN.CENTER)
    textbox(slide, x + 1.50, y + 0.20, 1.75, 0.34, title, 12.5, NAVY, True)
    textbox(slide, x + 1.18, y + 0.62, 2.10, 0.52, body, 8.4, TEXT)


def small_node(slide, x, y, title, icon):
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.12), Inches(y + 0.55), Inches(1.24), Inches(0.30))
    oval.fill.solid()
    oval.fill.fore_color.rgb = RGBColor(224, 241, 255)
    oval.line.color.rgb = RGBColor(210, 231, 252)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.28), Inches(y + 0.08), Inches(0.92), Inches(0.46))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(103, 173, 246)
    shape.line.width = Pt(1.0)
    set_text(shape, icon, 7.8, NAVY, True, PP_ALIGN.CENTER)
    textbox(slide, x - 0.02, y + 0.78, 1.52, 0.26, title, 8.4, TEXT, True, PP_ALIGN.CENTER)


def metric_block(slide, x, title, icon, body):
    textbox(slide, x, 6.46, 1.25, 0.23, title, 8.8, NAVY, True, PP_ALIGN.CENTER)
    textbox(slide, x + 0.34, 6.94, 0.52, 0.42, icon, 20, BLUE, True, PP_ALIGN.CENTER)
    textbox(slide, x + 0.08, 7.48, 1.09, 0.43, body, 6.2, TEXT, False, PP_ALIGN.CENTER)


def value_item(slide, x, icon, title, body):
    textbox(slide, x, 6.92, 0.55, 0.35, icon, 17, BLUE, True, PP_ALIGN.CENTER)
    textbox(slide, x - 0.16, 7.32, 0.88, 0.24, title, 6.8, TEXT, False, PP_ALIGN.CENTER)
    textbox(slide, x - 0.16, 7.58, 0.88, 0.24, body, 6.8, TEXT, False, PP_ALIGN.CENTER)


def build_deck(output_path: str | Path = "output/power_observability_editable.pptx") -> Path:
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    textbox(slide, 0.58, -0.08, 11.6, 0.52, "电力可观测性架构图可编辑重建示例", 25, NAVY, True)
    textbox(slide, 0.82, 0.42, 7.0, 0.34, "用 python-pptx 将架构图重建为可编辑 PowerPoint 原生对象", 14, TEXT)
    connector(slide, 0.50, 0.84, 15.62, 0.84, RGBColor(145, 190, 244), 1.0)

    risk_card(slide, 0.58, 1.18, 1, "链路复杂", "业务链路长，跨系统排障难", "1")
    risk_card(slide, 4.28, 1.18, 2, "证据分散", "故障复盘依赖人工拼接证据", "2")
    risk_card(slide, 7.98, 1.18, 3, "可视不足", "业务体验问题难以快速感知", "3")
    risk_card(slide, 11.68, 1.18, 4, "闭环偏弱", "问题处置与复盘缺少统一视图", "4")

    box(slide, 3.18, 3.08, 9.52, 2.62, LIGHT_BLUE, RGBColor(160, 204, 248), True, 1.3)
    textbox(slide, 6.18, 3.22, 3.52, 0.34, "可观测性能力平台", 18, NAVY, True, PP_ALIGN.CENTER)

    small_node(slide, 3.58, 3.88, "数据采集", "A")
    small_node(slide, 5.28, 3.88, "链路追踪", "B")
    small_node(slide, 6.98, 3.88, "异常识别", "C")
    small_node(slide, 8.68, 3.88, "根因定位", "D")
    small_node(slide, 10.38, 3.88, "证据回溯", "E")

    connector(slide, 4.35, 4.18, 5.28, 4.18)
    connector(slide, 6.05, 4.18, 6.98, 4.18)
    connector(slide, 7.75, 4.18, 8.68, 4.18)
    connector(slide, 9.45, 4.18, 10.38, 4.18)

    metric_block(slide, 2.22, "指标采集", "●", "分钟级")
    metric_block(slide, 4.08, "实时监测", "◆", "秒级")
    metric_block(slide, 5.94, "主动预警", "▲", "自动")
    metric_block(slide, 7.80, "高效排障", "■", "闭环")

    value_item(slide, 10.72, "◉", "业务体验", "可量化")
    value_item(slide, 11.82, "◎", "证据沉淀", "可追溯")
    value_item(slide, 12.92, "◇", "责任边界", "可定界")
    value_item(slide, 14.02, "⬢", "复盘优化", "可闭环")

    textbox(slide, 0.78, 6.22, 1.08, 0.32, "输入侧", 12, NAVY, True, PP_ALIGN.CENTER)
    box(slide, 0.62, 6.58, 1.38, 1.02, WHITE, RGBColor(158, 202, 246), True, 1.0)
    textbox(slide, 0.74, 6.82, 1.12, 0.52, "网络设备\n业务系统\n日志数据", 9.5, TEXT, False, PP_ALIGN.CENTER)

    connector(slide, 2.02, 7.08, 3.10, 7.08, BLUE2, 1.4)
    connector(slide, 12.72, 7.08, 14.72, 7.08, BLUE2, 1.4)

    box(slide, 14.80, 6.58, 1.44, 1.02, WHITE, RGBColor(158, 202, 246), True, 1.0)
    textbox(slide, 14.92, 6.82, 1.18, 0.52, "告警联动\n报告输出\n运维协同", 9.5, TEXT, False, PP_ALIGN.CENTER)
    textbox(slide, 14.90, 6.22, 1.22, 0.32, "输出侧", 12, NAVY, True, PP_ALIGN.CENTER)

    prs.save(output)
    return output


if __name__ == "__main__":
    out = build_deck()
    print(out)
